"""
文件解析模块 — 将 PDF/MD 转为统一的 content_list 格式。

职责单一：只负责「文件 → content_list」，不涉及切块或向量化。
可插拔：通过注册新的解析器支持更多文件类型。
"""

import json
import re
import time
from pathlib import Path
from typing import Any, Dict, List

from rag.config import ParserConfig


# ── 解析器注册表 ──────────────────────────────────────────

_PARSERS: Dict[str, type] = {}


def register_parser(extensions: List[str]):
    """装饰器：注册文件解析器。"""
    def decorator(cls):
        for ext in extensions:
            _PARSERS[ext.lower()] = cls
        return cls
    return decorator


def get_parser(file_path: Path, config: ParserConfig) -> "BaseParser":
    """根据文件扩展名获取对应解析器实例。"""
    ext = file_path.suffix.lower()
    parser_cls = _PARSERS.get(ext)
    if parser_cls is None:
        raise ValueError(f"不支持的文件类型: {ext}，已注册: {list(_PARSERS.keys())}")
    return parser_cls(config)


# ── 基类 ──────────────────────────────────────────────────

class BaseParser:
    """解析器基类。"""

    def __init__(self, config: ParserConfig):
        self.config = config

    def parse(self, file_path: Path) -> List[Dict[str, Any]]:
        """解析文件，返回 content_list。"""
        raise NotImplementedError


# ── PDF 解析器 ────────────────────────────────────────────

@register_parser([".pdf"])
class PDFParser(BaseParser):
    """使用 MinerU 解析 PDF。

    支持两种模式（通过 config.mineru_mode 控制）：
      - "local": 直接调用 do_parse()，需本地模型和 torch
      - "api":   通过 HTTP API 调用 MinerU 服务（服务独立部署）
    """

    def parse(self, file_path: Path) -> List[Dict[str, Any]]:
        if self.config.mineru_mode == "api":
            return self._parse_via_api(file_path)
        else:
            return self._parse_local(file_path)

    # ── 本地模式 ──────────────────────────────────────────

    def _parse_local(self, file_path: Path) -> List[Dict[str, Any]]:
        import os
        os.environ.setdefault("MINERU_MODEL_SOURCE", "local")
        os.environ.setdefault("HF_HUB_OFFLINE", "1")

        from mineru.cli.common import do_parse, read_fn

        stem = file_path.stem
        size_kb = file_path.stat().st_size / 1024
        print(f"[Parser] PDF (local): {file_path.name} ({size_kb:.0f} KB)")

        pdf_bytes = read_fn(file_path)

        print(f"[Parser] MinerU 解析中... ({self.config.backend}/{self.config.parse_method})")
        t0 = time.time()
        do_parse(
            output_dir=self.config.output_dir,
            pdf_file_names=[stem],
            pdf_bytes_list=[pdf_bytes],
            p_lang_list=[self.config.lang],
            backend=self.config.backend,
            parse_method=self.config.parse_method,
            formula_enable=self.config.formula_enable,
            table_enable=self.config.table_enable,
        )
        elapsed = time.time() - t0
        print(f"[Parser] 解析完成 ({elapsed:.1f}s)")

        # 读取 content_list
        content_path = Path(self.config.output_dir) / stem / "auto" / f"{stem}_content_list.json"
        if not content_path.exists():
            raise FileNotFoundError(f"MinerU 未生成: {content_path}")

        content_list = json.loads(content_path.read_text(encoding="utf-8"))
        print(f"[Parser] 输出 {len(content_list)} 个 block")
        return content_list

    # ── API 模式（MinerU 云端 API）──────────────────────

    _MINERU_API_BASE = "https://mineru.net/api/v4"

    def _parse_via_api(self, file_path: Path) -> List[Dict[str, Any]]:
        import re
        import tempfile
        import zipfile
        from pathlib import PurePosixPath

        import httpx

        token = self.config.mineru_api_token
        if not token:
            raise ValueError(
                "MinerU API 模式需要 token，请在 .env 中设置 RAG_MINERU_API_TOKEN\n"
                "获取地址: https://mineru.net/apiManage/token"
            )

        stem = file_path.stem
        size_kb = file_path.stat().st_size / 1024
        print(f"[Parser] PDF (MinerU Cloud API): {file_path.name} ({size_kb:.0f} KB)")
        print(f"[Parser] 模型版本: {self.config.mineru_model_version}")

        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {token}",
        }
        base = self._MINERU_API_BASE
        t0 = time.time()

        # ── Step 1: 申请预签名上传 URL ──
        print(f"[Parser] Step 1/4: 申请上传链接...")
        body = {
            "files": [{"name": file_path.name}],
            "model_version": self.config.mineru_model_version,
            "enable_formula": self.config.formula_enable,
            "enable_table": self.config.table_enable,
            "language": self.config.lang,
        }
        resp = httpx.post(f"{base}/file-urls/batch", headers=headers, json=body, timeout=30)
        if resp.status_code != 200:
            self._raise_api_error("申请上传链接失败", resp)
        result = resp.json()
        if result.get("code") != 0:
            raise RuntimeError(f"申请上传链接失败: {result.get('msg', result)}")

        batch_id = result["data"]["batch_id"]
        file_urls = result["data"]["file_urls"]
        print(f"[Parser] batch_id={batch_id}, 获得 {len(file_urls)} 个上传链接")

        # ── Step 2: PUT 上传文件 ──
        print(f"[Parser] Step 2/4: 上传文件...")
        file_bytes = file_path.read_bytes()
        for i, upload_url in enumerate(file_urls):
            put_resp = httpx.put(upload_url, content=file_bytes, timeout=120)
            if put_resp.status_code != 200:
                raise RuntimeError(f"文件上传失败 ({put_resp.status_code}): {put_resp.text[:200]}")
        print(f"[Parser] 上传完成 ({len(file_bytes)/1024:.0f} KB)")

        # ── Step 3: 轮询解析结果 ──
        print(f"[Parser] Step 3/4: 等待解析完成...")
        poll_url = f"{base}/extract-results/batch/{batch_id}"
        max_wait = 600
        interval = 3.0
        waited = 0.0

        while waited < max_wait:
            import time as _time
            _time.sleep(interval)
            waited += interval

            resp = httpx.get(poll_url, headers=headers, timeout=30)
            if resp.status_code != 200:
                self._raise_api_error("查询任务状态失败", resp)
            result = resp.json()
            if result.get("code") != 0:
                raise RuntimeError(f"查询失败: {result.get('msg')}")

            extract_results = result["data"].get("extract_result", [])
            if not extract_results:
                continue

            state = extract_results[0].get("state")
            if state == "done":
                full_zip_url = extract_results[0].get("full_zip_url")
                if not full_zip_url:
                    raise RuntimeError("解析完成但未返回 full_zip_url")
                break
            elif state == "failed":
                err = extract_results[0].get("err_msg", "未知错误")
                raise RuntimeError(f"解析失败: {err}")
            else:
                progress = extract_results[0].get("extract_progress", {})
                if progress:
                    print(f"[Parser] 进度: {progress.get('extracted_pages','?')}/"
                          f"{progress.get('total_pages','?')} 页, 状态: {state}")

        if waited >= max_wait:
            raise RuntimeError(f"解析超时 ({max_wait}s)")

        elapsed = time.time() - t0
        print(f"[Parser] 解析完成 ({elapsed:.1f}s)")

        # ── Step 4: 下载并解压结果 ──
        print(f"[Parser] Step 4/4: 下载结果...")
        resp = httpx.get(full_zip_url, timeout=60)
        if resp.status_code != 200:
            raise RuntimeError(f"下载结果失败: {resp.status_code}")

        _M = PurePosixPath  # Use PurePosixPath for ZIP internal paths
        with tempfile.NamedTemporaryFile(suffix=".zip", delete=False) as tmp:
            tmp.write(resp.content)
            zip_path = tmp.name

        try:
            with zipfile.ZipFile(zip_path, "r") as zf:
                all_names = zf.namelist()
                print(f"[Parser] ZIP 内容: {len(all_names)} 个文件")

                # 优先查找结构化 JSON（layout/*.json），其次用 markdown 转 content_list
                layout_jsons = sorted(
                    [n for n in all_names if n.endswith(".json") and "layout" in n.lower()],
                )
                if layout_jsons:
                    # 有 layout JSON，合并所有页的结构信息
                    content_list = self._layout_json_to_content_list(zf, layout_jsons)
                else:
                    # 回退：用 markdown 转 content_list
                    md_files = [n for n in all_names if n.endswith(".md")]
                    if not md_files:
                        raise FileNotFoundError(
                            f"ZIP 中未找到可解析的文件。内容: {all_names}"
                        )
                    content_list = self._markdown_zip_to_content_list(zf, md_files)

                print(f"[Parser] 输出 {len(content_list)} 个 block")
                return content_list
        finally:
            Path(zip_path).unlink(missing_ok=True)

    def _layout_json_to_content_list(
        self, zf, layout_paths: list
    ) -> List[Dict[str, Any]]:
        """将 MinerU 云端 API 返回的 layout/*.json 转为 content_list 格式。"""
        items = []
        for path in layout_paths:
            data = json.loads(zf.read(path).decode("utf-8"))
            # layout.json 结构通常是 list，每项包含 type/text/bbox/page_idx 等
            if isinstance(data, list):
                for block in data:
                    item = {"type": block.get("type", "text")}
                    if "text" in block:
                        item["text"] = block["text"]
                    elif "content" in block:
                        item["text"] = block["content"]
                    if "page_idx" in block or "page_id" in block:
                        item["page_idx"] = block.get("page_idx", block.get("page_id", 0))
                    if "bbox" in block:
                        item["bbox"] = block["bbox"]
                    if "level" in block:
                        item["text_level"] = block["level"]
                    if item.get("text"):
                        items.append(item)
            elif isinstance(data, dict):
                for page_blocks in data.values():
                    if isinstance(page_blocks, list):
                        for block in page_blocks:
                            if isinstance(block, dict) and block.get("text"):
                                items.append({
                                    "type": block.get("type", "text"),
                                    "text": block.get("text", block.get("content", "")),
                                    "page_idx": block.get("page_idx", block.get("page_id", 0)),
                                })
        if not items:
            raise RuntimeError("layout JSON 中无有效内容")
        return items

    def _markdown_zip_to_content_list(
        self, zf, md_paths: list
    ) -> List[Dict[str, Any]]:
        """将 MinerU 云端 API 返回的 markdown 文件转为 content_list 格式。"""
        # 复用 MarkdownParser 的分段逻辑
        from rag.parser import MarkdownParser, ParserConfig as _PC
        import tempfile
        import os

        # 合并所有 md 文件
        full_md = ""
        for path in sorted(md_paths):
            full_md += zf.read(path).decode("utf-8") + "\n\n"

        # 使用 MarkdownParser 的内部方法
        mp = MarkdownParser.__new__(MarkdownParser)
        mp.config = self.config
        return mp._md_to_content_list(full_md)

    @staticmethod
    def _raise_api_error(prefix: str, resp) -> None:
        detail = ""
        try:
            detail = resp.json().get("msg", resp.text[:300])
        except Exception:
            detail = resp.text[:300]
        raise RuntimeError(f"{prefix} ({resp.status_code}): {detail}")


# ── PPTX 解析器 ───────────────────────────────────────────

@register_parser([".ppt"])
class PPTParser(BaseParser):
    """解析旧版 .ppt（二进制格式），通过 COM/ LibreOffice 转为 .pptx 后解析。"""

    def parse(self, file_path: Path) -> List[Dict[str, Any]]:
        import subprocess
        import tempfile
        import os
        import platform

        size_kb = file_path.stat().st_size / 1024
        print(f"[Parser] PPT (legacy): {file_path.name} ({size_kb:.0f} KB)")

        # 创建临时 .pptx 文件
        tmp_dir = tempfile.mkdtemp(prefix="ppt_convert_")
        pptx_path = Path(tmp_dir) / f"{file_path.stem}.pptx"

        try:
            if platform.system() == "Windows":
                self._convert_via_com(file_path, pptx_path)
            else:
                self._convert_via_libreoffice(file_path, pptx_path)

            # 委托 PPTXParser 解析
            pptx_parser = PPTXParser(self.config)
            return pptx_parser.parse(pptx_path)

        finally:
            # 清理临时文件
            import shutil
            try:
                shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception:
                pass

    def _convert_via_com(self, source: Path, target: Path):
        """Windows: 通过 PowerPoint COM 自动化转换。"""
        import subprocess

        ps1 = Path(__file__).parent / "convert_ppt.ps1"
        result = subprocess.run(
            [
                "powershell.exe", "-NoProfile", "-ExecutionPolicy", "Bypass",
                "-File", str(ps1),
                "-SourceFile", str(source.resolve()),
                "-TargetFile", str(target.resolve()),
            ],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0 or not target.exists():
            err = result.stderr.strip() or result.stdout.strip()
            raise RuntimeError(
                f"PPT 转换失败（需要安装 PowerPoint）。\n"
                f"请手动将 .ppt 另存为 .pptx 格式后再上传。\n"
                f"错误详情: {err[:500]}"
            )
        print(f"[Parser] PPT → PPTX 转换成功")

    def _convert_via_libreoffice(self, source: Path, target: Path):
        """非 Windows: 通过 LibreOffice CLI 转换。"""
        import subprocess

        src = str(source.resolve())
        outdir = str(target.parent.resolve())
        result = subprocess.run(
            [
                "libreoffice", "--headless", "--convert-to", "pptx",
                "--outdir", outdir, src,
            ],
            capture_output=True, text=True, timeout=120,
        )
        # LibreOffice names output as <stem>.pptx automatically
        auto_name = target.parent / f"{source.stem}.pptx"
        if result.returncode != 0 or not auto_name.exists():
            err = result.stderr.strip() or result.stdout.strip()
            raise RuntimeError(
                f"PPT 转换失败（需要安装 LibreOffice）。\n"
                f"请手动将 .ppt 另存为 .pptx 格式后再上传。\n"
                f"错误详情: {err[:500]}"
            )
        # Rename to expected target path
        if auto_name != target:
            import shutil
            shutil.move(str(auto_name), str(target))
        print(f"[Parser] PPT → PPTX 转换成功 (LibreOffice)")


@register_parser([".pptx"])
class PPTXParser(BaseParser):
    """使用 python-pptx 解析 PowerPoint 文件。"""

    def parse(self, file_path: Path) -> List[Dict[str, Any]]:
        from pptx import Presentation

        size_kb = file_path.stat().st_size / 1024
        print(f"[Parser] PPTX: {file_path.name} ({size_kb:.0f} KB)")

        prs = Presentation(str(file_path))
        items: List[Dict[str, Any]] = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    table = shape.table
                    rows_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        rows_data.append(" | ".join(row_data))
                    table_text = "\n".join(rows_data)
                    if table_text.strip():
                        items.append({
                            "type": "table",
                            "text": table_text,
                            "page_idx": slide_idx,
                            "bbox": [0, 0, 0, 0],
                        })

            if slide_texts:
                # 第一行作为标题
                items.append({
                    "type": "text",
                    "text": slide_texts[0],
                    "text_level": 2,
                    "page_idx": slide_idx,
                    "bbox": [0, 0, 0, 0],
                })
                if len(slide_texts) > 1:
                    items.append({
                        "type": "text",
                        "text": "\n".join(slide_texts[1:]),
                        "page_idx": slide_idx,
                        "bbox": [0, 0, 0, 0],
                    })

        print(f"[Parser] 输出 {len(items)} 个 block")
        return items


# ── DOCX 解析器 ───────────────────────────────────────────

@register_parser([".docx"])
class DOCXParser(BaseParser):
    """使用 python-docx 解析 Word 文件。"""

    def parse(self, file_path: Path) -> List[Dict[str, Any]]:
        from docx import Document

        size_kb = file_path.stat().st_size / 1024
        print(f"[Parser] DOCX: {file_path.name} ({size_kb:.0f} KB)")

        doc = Document(str(file_path))
        items: List[Dict[str, Any]] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # 根据段落样式判断标题级别
            style_name = (para.style.name or "").lower()
            if "heading" in style_name:
                # 提取标题级别 (Heading 1 → 1, Heading 2 → 2, ...)
                level = 1
                for ch in style_name:
                    if ch.isdigit():
                        level = int(ch)
                        break
                items.append({
                    "type": "text",
                    "text": text,
                    "text_level": level,
                    "page_idx": 0,
                    "bbox": [0, 0, 0, 0],
                })
            else:
                items.append({
                    "type": "text",
                    "text": text,
                    "page_idx": 0,
                    "bbox": [0, 0, 0, 0],
                })

        # 解析表格
        for table in doc.tables:
            rows_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                rows_data.append(" | ".join(row_data))
            table_text = "\n".join(rows_data)
            if table_text.strip():
                items.append({
                    "type": "table",
                    "text": table_text,
                    "page_idx": 0,
                    "bbox": [0, 0, 0, 0],
                })

        print(f"[Parser] 输出 {len(items)} 个 block")
        return items


# ── Excel 解析器 ──────────────────────────────────────────

@register_parser([".xlsx", ".xls"])
class ExcelParser(BaseParser):
    """使用 openpyxl 解析 Excel 文件。"""

    def parse(self, file_path: Path) -> List[Dict[str, Any]]:
        from openpyxl import load_workbook

        size_kb = file_path.stat().st_size / 1024
        print(f"[Parser] Excel: {file_path.name} ({size_kb:.0f} KB)")

        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        items: List[Dict[str, Any]] = []

        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]

            # Sheet 名作为标题
            items.append({
                "type": "text",
                "text": f"工作表: {sheet_name}",
                "text_level": 2,
                "page_idx": sheet_idx,
                "bbox": [0, 0, 0, 0],
            })

            # 将表格内容转为文本
            rows_data = []
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_text):
                    rows_data.append(" | ".join(row_text))

            if rows_data:
                items.append({
                    "type": "table",
                    "text": "\n".join(rows_data),
                    "page_idx": sheet_idx,
                    "bbox": [0, 0, 0, 0],
                })

        wb.close()
        print(f"[Parser] 输出 {len(items)} 个 block")
        return items


# ── Markdown 解析器 ───────────────────────────────────────

_HEADING_RE = re.compile(r'^(#{1,6})\s+(.+)$', re.MULTILINE)


@register_parser([".md", ".markdown", ".txt"])
class MarkdownParser(BaseParser):
    """将 Markdown/文本文件转为 content_list 格式。"""

    def parse(self, file_path: Path) -> List[Dict[str, Any]]:
        size_kb = file_path.stat().st_size / 1024
        print(f"[Parser] MD: {file_path.name} ({size_kb:.0f} KB)")

        md_text = self._read_file(file_path)
        content_list = self._md_to_content_list(md_text)
        print(f"[Parser] 转换为 {len(content_list)} 个 block")
        return content_list

    def _read_file(self, path: Path) -> str:
        """尝试多种编码读取文件。"""
        for enc in ('utf-8', 'gbk', 'gb2312', 'latin-1'):
            try:
                return path.read_text(encoding=enc)
            except UnicodeDecodeError:
                continue
        raise ValueError(f"无法读取文件: {path}")

    def _md_to_content_list(self, md_text: str) -> List[Dict[str, Any]]:
        """将 markdown 文本转为 content_list 格式。"""
        items: List[Dict[str, Any]] = []
        blocks = re.split(r'\n\s*\n', md_text.strip())

        for block in blocks:
            lines = block.strip().split('\n')
            if not lines:
                continue

            first = lines[0].strip()
            m = _HEADING_RE.match(first)
            if m:
                level = len(m.group(1))
                title_text = m.group(2).strip()
                items.append({
                    "type": "text",
                    "text": title_text,
                    "text_level": level,
                    "page_idx": 0,
                    "bbox": [0, 0, 0, 0],
                })
                # 标题行之后的剩余内容作为正文
                rest = '\n'.join(lines[1:]).strip()
                if rest:
                    items.append({
                        "type": "text",
                        "text": rest,
                        "page_idx": 0,
                        "bbox": [0, 0, 0, 0],
                    })
            else:
                items.append({
                    "type": "text",
                    "text": block.strip(),
                    "page_idx": 0,
                    "bbox": [0, 0, 0, 0],
                })

        return items
